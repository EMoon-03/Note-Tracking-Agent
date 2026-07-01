"""
file_processor.py — Reads any supported file and prepares its content for the AI

CONCEPT: Ollama (our local AI) accepts text and images as input. This module
figures out what kind of file you gave it and extracts the content in the right
format. For example:
  - Text files → just read the text
  - Images → pass the file path directly (Ollama can read image files)
  - PDFs → extract the text using pypdf (Ollama can't read raw PDF binary)
  - PowerPoints → extract slide text using python-pptx
"""

from pathlib import Path
from pypdf import PdfReader
from pptx import Presentation


# ─── Supported file types ─────────────────────────────────────────────────────

SUPPORTED_EXTENSIONS = {
    ".txt": "text",
    ".md": "text",
    ".csv": "text",
    ".py": "text",
    ".js": "text",
    ".html": "text",
    ".jpg": "image",
    ".jpeg": "image",
    ".png": "image",
    ".gif": "image",
    ".webp": "image",
    ".pdf": "pdf",
    ".pptx": "pptx",
}


def process_file(file_path: str) -> dict:
    """
    Read a file and return its content ready to send to Ollama.

    Returns a dict with:
        {
            "text": "the text content or prompt",
            "images": []           # empty list, OR
            "images": ["/path"]    # file path for images
        }

    Ollama accepts both text and image inputs in a single message.
    """
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = path.suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(SUPPORTED_EXTENSIONS.keys())
        raise ValueError(f"Unsupported file type: '{ext}'\nSupported: {supported}")

    file_type = SUPPORTED_EXTENSIONS[ext]

    if file_type == "text":
        return _process_text(path)
    elif file_type == "image":
        return _process_image(path)
    elif file_type == "pdf":
        return _process_pdf(path)
    elif file_type == "pptx":
        return _process_pptx(path)


# ─── Private helpers ──────────────────────────────────────────────────────────


def _process_text(path: Path) -> dict:
    text = path.read_text(encoding="utf-8")
    return {
        "text": f"Here is the content of '{path.name}':\n\n{text}",
        "images": [],
    }


def _process_image(path: Path) -> dict:
    """
    For images, we just pass the file path to Ollama.
    Ollama's vision models (like llava) can read image files directly —
    no need to convert to base64 like we'd need with a remote API.
    """
    return {
        "text": f"The image is named '{path.name}'. Generate notes from its content.",
        "images": [str(path.resolve())],  # absolute path
    }


def _process_pdf(path: Path) -> dict:
    """
    Extract text from a PDF using pypdf.

    CONCEPT: Unlike Claude (which understands PDFs natively), local models
    like llava work with plain text. We use pypdf to pull the text out of
    each page, then send that text to the model.
    """
    reader = PdfReader(str(path))

    pages_text = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text()
        if text and text.strip():
            pages_text.append(f"--- Page {i} ---\n{text.strip()}")

    if not pages_text:
        raise ValueError(
            "No text could be extracted from this PDF. "
            "It may be a scanned image-only PDF."
        )

    full_text = "\n\n".join(pages_text)
    return {
        "text": f"Here is the text extracted from '{path.name}':\n\n{full_text}",
        "images": [],
    }


def _process_pptx(path: Path) -> dict:
    """Extract slide text from a PowerPoint file using python-pptx."""
    prs = Presentation(str(path))

    slides_text = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if slide_lines:
            slides_text.append(f"### Slide {i}\n" + "\n".join(slide_lines))

    if not slides_text:
        raise ValueError("No text content found in this PowerPoint file.")

    return {
        "text": (
            f"Here is the text from the PowerPoint '{path.name}':\n\n"
            + "\n\n".join(slides_text)
        ),
        "images": [],
    }
